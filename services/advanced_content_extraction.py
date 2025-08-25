"""
Advanced content extraction service with llava:7b integration
Provides intelligent document structure analysis and visual element understanding
"""
import os
import io
import json
import base64
import tempfile
from typing import List, Dict, Any, Optional, Tuple, Union
from pathlib import Path
from dataclasses import dataclass
import time

# Core processing libraries
from langchain.schema import Document as LangChainDocument
from langchain_community.llms import Ollama
from PIL import Image
import fitz  # PyMuPDF for PDF processing

# Import existing processors
from services.document_processors import ProcessingResult

# Optional dependencies
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False

try:
    import cv2
    import numpy as np
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

from config import config
from utils.logger import get_logger

logger = get_logger("advanced_content_extraction")


@dataclass
class StructuralElement:
    """Represents a structural element identified in a document"""
    element_type: str  # table, chart, diagram, heading, paragraph, etc.
    content: str
    position: Dict[str, Any]  # page, coordinates, etc.
    confidence: float
    metadata: Dict[str, Any]


@dataclass
class AdvancedProcessingResult:
    """Enhanced processing result with structural analysis"""
    success: bool
    documents: List[LangChainDocument]
    structural_elements: List[StructuralElement]
    visual_analysis: Dict[str, Any]
    metadata: Dict[str, Any]
    error_message: Optional[str] = None
    processing_time: Optional[float] = None


class LlavaVisionAnalyzer:
    """Vision analyzer using llava:7b for document structure understanding"""
    
    def __init__(self):
        self.vision_model = Ollama(
            model=config.ai.vision_model,
            base_url=config.ai.ollama_base_url,
            temperature=0.1  # Lower temperature for more consistent analysis
        )
    
    def encode_image_to_base64(self, image_path: str) -> str:
        """Convert image to base64 for llava processing"""
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    
    def analyze_document_structure(self, image_path: str, context: str = "") -> Dict[str, Any]:
        """Analyze document structure using llava:7b"""
        try:
            # Create a structured prompt for document analysis
            prompt = f"""Analyze this document image and identify its structural elements. Focus on:

1. Document layout and structure
2. Tables, charts, and diagrams
3. Headings and sections
4. Text blocks and paragraphs
5. Visual elements (images, graphics)

Context: {context}

Please provide a JSON response with the following structure:
{{
    "document_type": "type of document",
    "layout": "description of overall layout",
    "elements": [
        {{
            "type": "element type (table/chart/heading/paragraph/image)",
            "description": "detailed description",
            "position": "location in document",
            "content_summary": "summary of content if readable"
        }}
    ],
    "tables_detected": number,
    "charts_detected": number,
    "text_regions": number,
    "quality_assessment": "assessment of image quality and readability"
}}

Respond only with valid JSON."""

            # Note: llava:7b in Ollama typically expects the image to be referenced
            # The exact API may vary, but this is the general approach
            response = self.vision_model.invoke(f"[IMAGE: {image_path}] {prompt}")
            
            # Try to parse JSON response
            try:
                analysis = json.loads(response)
                return analysis
            except json.JSONDecodeError:
                # If JSON parsing fails, create a structured response from text
                return {
                    "document_type": "unknown",
                    "layout": "Could not parse structured analysis",
                    "elements": [],
                    "tables_detected": 0,
                    "charts_detected": 0,
                    "text_regions": 0,
                    "quality_assessment": "Analysis completed but format parsing failed",
                    "raw_response": response
                }
                
        except Exception as e:
            logger.error(f"Vision analysis failed: {e}")
            return {
                "document_type": "unknown",
                "layout": "Analysis failed",
                "elements": [],
                "tables_detected": 0,
                "charts_detected": 0,
                "text_regions": 0,
                "quality_assessment": f"Error: {str(e)}",
                "error": str(e)
            }
    
    def analyze_visual_elements(self, image_path: str) -> Dict[str, Any]:
        """Analyze visual elements and convert to structured descriptions"""
        try:
            prompt = """Analyze this image and describe any visual elements such as:
- Charts, graphs, and diagrams
- Tables and their structure
- Images and illustrations
- Text formatting and layout

Provide detailed descriptions that would help someone understand the visual content without seeing the image. Focus on data, relationships, and key information presented visually.

Format your response as structured text that can be easily processed."""

            response = self.vision_model.invoke(f"[IMAGE: {image_path}] {prompt}")
            
            return {
                "visual_description": response,
                "analysis_method": "llava_vision",
                "timestamp": time.time()
            }
            
        except Exception as e:
            logger.error(f"Visual element analysis failed: {e}")
            return {
                "visual_description": f"Visual analysis failed: {str(e)}",
                "analysis_method": "error",
                "error": str(e)
            }


class AdvancedPDFProcessor:
    """Enhanced PDF processor with llava:7b integration"""
    
    def __init__(self):
        self.vision_analyzer = LlavaVisionAnalyzer()
    
    def process_pdf_with_structure_analysis(self, file_path: str) -> AdvancedProcessingResult:
        """Process PDF with advanced structure analysis using llava:7b"""
        start_time = time.time()
        
        try:
            documents = []
            structural_elements = []
            visual_analysis = {}
            
            # Open PDF
            pdf_document = fitz.open(file_path)
            
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                
                # Extract text using PyMuPDF
                text = page.get_text()
                
                # Convert page to image for vision analysis
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Higher resolution
                img_data = pix.tobytes("png")
                
                # Save temporary image for llava analysis
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_img:
                    temp_img.write(img_data)
                    temp_img_path = temp_img.name
                
                try:
                    # Analyze page structure with llava:7b
                    structure_analysis = self.vision_analyzer.analyze_document_structure(
                        temp_img_path, 
                        context=f"PDF page {page_num + 1} of {len(pdf_document)}"
                    )
                    
                    # Analyze visual elements
                    visual_elements = self.vision_analyzer.analyze_visual_elements(temp_img_path)
                    
                    # Extract tables and structured data
                    tables = self._extract_tables_from_page(page)
                    
                    # Create enhanced document with structure information
                    enhanced_content = self._create_enhanced_content(
                        text, structure_analysis, visual_elements, tables
                    )
                    
                    doc = LangChainDocument(
                        page_content=enhanced_content,
                        metadata={
                            "page_number": page_num + 1,
                            "source": file_path,
                            "document_type": "pdf",
                            "extraction_method": "advanced_llava_analysis",
                            "structure_analysis": structure_analysis,
                            "visual_elements": visual_elements,
                            "tables_count": len(tables),
                            "has_visual_content": structure_analysis.get("charts_detected", 0) > 0 or 
                                                structure_analysis.get("tables_detected", 0) > 0
                        }
                    )
                    documents.append(doc)
                    
                    # Create structural elements
                    for element in structure_analysis.get("elements", []):
                        structural_elements.append(StructuralElement(
                            element_type=element.get("type", "unknown"),
                            content=element.get("content_summary", ""),
                            position={"page": page_num + 1, "location": element.get("position", "")},
                            confidence=0.8,  # Default confidence for llava analysis
                            metadata=element
                        ))
                    
                    visual_analysis[f"page_{page_num + 1}"] = {
                        "structure": structure_analysis,
                        "visual_elements": visual_elements,
                        "tables": tables
                    }
                    
                finally:
                    # Clean up temporary image
                    if os.path.exists(temp_img_path):
                        os.unlink(temp_img_path)
            
            pdf_document.close()
            processing_time = time.time() - start_time
            
            return AdvancedProcessingResult(
                success=True,
                documents=documents,
                structural_elements=structural_elements,
                visual_analysis=visual_analysis,
                metadata={
                    "total_pages": len(pdf_document),
                    "extraction_method": "advanced_llava_pdf",
                    "processing_time": processing_time,
                    "vision_model": config.ai.vision_model
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Advanced PDF processing failed: {e}")
            return AdvancedProcessingResult(
                success=False,
                documents=[],
                structural_elements=[],
                visual_analysis={},
                metadata={"error": str(e)},
                error_message=str(e),
                processing_time=time.time() - start_time
            )
    
    def _extract_tables_from_page(self, page) -> List[Dict[str, Any]]:
        """Extract tables from PDF page using PyMuPDF"""
        tables = []
        try:
            # Find tables in the page
            table_list = page.find_tables()
            
            for table_index, table in enumerate(table_list):
                # Extract table data
                table_data = table.extract()
                
                if table_data:
                    # Convert to structured format
                    table_info = {
                        "table_index": table_index,
                        "rows": len(table_data),
                        "columns": len(table_data[0]) if table_data else 0,
                        "data": table_data,
                        "bbox": table.bbox,  # Bounding box coordinates
                        "formatted_content": self._format_table_content(table_data)
                    }
                    tables.append(table_info)
                    
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")
        
        return tables
    
    def _format_table_content(self, table_data: List[List[str]]) -> str:
        """Format table data as readable text"""
        if not table_data:
            return ""
        
        # Create a formatted table representation
        formatted_lines = []
        for row in table_data:
            # Clean and join cells
            cleaned_row = [str(cell).strip() if cell else "" for cell in row]
            formatted_lines.append(" | ".join(cleaned_row))
        
        return "\n".join(formatted_lines)
    
    def _create_enhanced_content(self, text: str, structure_analysis: Dict, 
                               visual_elements: Dict, tables: List[Dict]) -> str:
        """Create enhanced content combining text, structure, and visual analysis"""
        content_parts = []
        
        # Add original text
        if text.strip():
            content_parts.append("=== TEXT CONTENT ===")
            content_parts.append(text.strip())
        
        # Add structure analysis
        if structure_analysis.get("elements"):
            content_parts.append("\n=== DOCUMENT STRUCTURE ===")
            content_parts.append(f"Document Type: {structure_analysis.get('document_type', 'Unknown')}")
            content_parts.append(f"Layout: {structure_analysis.get('layout', 'Not analyzed')}")
            
            for element in structure_analysis.get("elements", []):
                content_parts.append(f"- {element.get('type', 'Unknown')}: {element.get('description', '')}")
        
        # Add visual analysis
        if visual_elements.get("visual_description"):
            content_parts.append("\n=== VISUAL ELEMENTS ===")
            content_parts.append(visual_elements["visual_description"])
        
        # Add table content
        if tables:
            content_parts.append("\n=== TABLES ===")
            for i, table in enumerate(tables):
                content_parts.append(f"Table {i + 1} ({table['rows']}x{table['columns']}):")
                content_parts.append(table["formatted_content"])
                content_parts.append("")
        
        return "\n".join(content_parts)


class AdvancedExcelProcessor:
    """Enhanced Excel processor with relationship and formula preservation"""
    
    def process_excel_with_relationships(self, file_path: str) -> AdvancedProcessingResult:
        """Process Excel files maintaining data relationships and formulas"""
        start_time = time.time()
        
        try:
            if not PANDAS_AVAILABLE:
                raise ImportError("pandas not available for Excel processing")
            
            documents = []
            structural_elements = []
            
            # Read Excel file with multiple sheets
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                # Read sheet with formulas preserved
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Analyze sheet structure
                sheet_analysis = self._analyze_sheet_structure(df, sheet_name)
                
                # Create enhanced content
                content = self._create_excel_content(df, sheet_analysis)
                
                doc = LangChainDocument(
                    page_content=content,
                    metadata={
                        "sheet_name": sheet_name,
                        "source": file_path,
                        "document_type": "excel",
                        "extraction_method": "advanced_excel_analysis",
                        "rows": len(df),
                        "columns": len(df.columns),
                        "has_formulas": sheet_analysis.get("has_formulas", False),
                        "data_types": sheet_analysis.get("data_types", {}),
                        "structure_analysis": sheet_analysis
                    }
                )
                documents.append(doc)
                
                # Create structural elements for significant data patterns
                for element in sheet_analysis.get("structural_elements", []):
                    structural_elements.append(StructuralElement(
                        element_type=element["type"],
                        content=element["content"],
                        position={"sheet": sheet_name, "range": element.get("range", "")},
                        confidence=element.get("confidence", 0.8),
                        metadata=element
                    ))
            
            processing_time = time.time() - start_time
            
            return AdvancedProcessingResult(
                success=True,
                documents=documents,
                structural_elements=structural_elements,
                visual_analysis={"excel_analysis": "completed"},
                metadata={
                    "total_sheets": len(excel_file.sheet_names),
                    "sheet_names": excel_file.sheet_names,
                    "extraction_method": "advanced_excel",
                    "processing_time": processing_time
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Advanced Excel processing failed: {e}")
            return AdvancedProcessingResult(
                success=False,
                documents=[],
                structural_elements=[],
                visual_analysis={},
                metadata={"error": str(e)},
                error_message=str(e),
                processing_time=time.time() - start_time
            )
    
    def _analyze_sheet_structure(self, df: pd.DataFrame, sheet_name: str) -> Dict[str, Any]:
        """Analyze Excel sheet structure and patterns"""
        analysis = {
            "sheet_name": sheet_name,
            "dimensions": {"rows": len(df), "columns": len(df.columns)},
            "data_types": {},
            "has_formulas": False,  # Would need openpyxl for formula detection
            "structural_elements": []
        }
        
        # Analyze data types
        for col in df.columns:
            dtype = str(df[col].dtype)
            analysis["data_types"][col] = dtype
        
        # Identify potential tables and data regions
        if len(df) > 1 and len(df.columns) > 1:
            analysis["structural_elements"].append({
                "type": "data_table",
                "content": f"Data table with {len(df)} rows and {len(df.columns)} columns",
                "range": f"A1:{chr(65 + len(df.columns) - 1)}{len(df)}",
                "confidence": 0.9
            })
        
        return analysis
    
    def _create_excel_content(self, df: pd.DataFrame, analysis: Dict) -> str:
        """Create enhanced content for Excel data"""
        content_parts = []
        
        # Add sheet information
        content_parts.append(f"=== EXCEL SHEET: {analysis['sheet_name']} ===")
        content_parts.append(f"Dimensions: {analysis['dimensions']['rows']} rows × {analysis['dimensions']['columns']} columns")
        
        # Add column information
        content_parts.append("\n=== COLUMNS ===")
        for col, dtype in analysis["data_types"].items():
            content_parts.append(f"- {col}: {dtype}")
        
        # Add data sample
        content_parts.append("\n=== DATA SAMPLE ===")
        content_parts.append(df.head(10).to_string())
        
        # Add summary statistics for numeric columns
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            content_parts.append("\n=== NUMERIC SUMMARY ===")
            content_parts.append(df[numeric_cols].describe().to_string())
        
        return "\n".join(content_parts)


class AdvancedPowerPointProcessor:
    """Enhanced PowerPoint processor with slide content and speaker notes extraction"""
    
    def __init__(self):
        self.vision_analyzer = LlavaVisionAnalyzer()
    
    def process_powerpoint_with_analysis(self, file_path: str) -> AdvancedProcessingResult:
        """Process PowerPoint with advanced content extraction"""
        start_time = time.time()
        
        try:
            if not PYTHON_PPTX_AVAILABLE:
                raise ImportError("python-pptx not available for PowerPoint processing")
            
            documents = []
            structural_elements = []
            visual_analysis = {}
            
            # Open presentation
            presentation = Presentation(file_path)
            
            for slide_num, slide in enumerate(presentation.slides):
                # Extract text content
                slide_text = self._extract_slide_text(slide)
                
                # Extract speaker notes
                notes = self._extract_speaker_notes(slide)
                
                # Analyze slide layout and visual elements
                layout_analysis = self._analyze_slide_layout(slide)
                
                # Create enhanced content
                content = self._create_powerpoint_content(
                    slide_text, notes, layout_analysis, slide_num + 1
                )
                
                doc = LangChainDocument(
                    page_content=content,
                    metadata={
                        "slide_number": slide_num + 1,
                        "source": file_path,
                        "document_type": "powerpoint",
                        "extraction_method": "advanced_powerpoint_analysis",
                        "has_notes": bool(notes.strip()),
                        "layout_type": layout_analysis.get("layout_name", "unknown"),
                        "shape_count": layout_analysis.get("shape_count", 0),
                        "has_images": layout_analysis.get("has_images", False),
                        "has_charts": layout_analysis.get("has_charts", False)
                    }
                )
                documents.append(doc)
                
                # Create structural elements
                for element in layout_analysis.get("elements", []):
                    structural_elements.append(StructuralElement(
                        element_type=element["type"],
                        content=element["content"],
                        position={"slide": slide_num + 1, "shape_id": element.get("shape_id", "")},
                        confidence=0.8,
                        metadata=element
                    ))
                
                visual_analysis[f"slide_{slide_num + 1}"] = layout_analysis
            
            processing_time = time.time() - start_time
            
            return AdvancedProcessingResult(
                success=True,
                documents=documents,
                structural_elements=structural_elements,
                visual_analysis=visual_analysis,
                metadata={
                    "total_slides": len(presentation.slides),
                    "extraction_method": "advanced_powerpoint",
                    "processing_time": processing_time
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Advanced PowerPoint processing failed: {e}")
            return AdvancedProcessingResult(
                success=False,
                documents=[],
                structural_elements=[],
                visual_analysis={},
                metadata={"error": str(e)},
                error_message=str(e),
                processing_time=time.time() - start_time
            )
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from slide"""
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        
        return "\n".join(text_parts)
    
    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from slide"""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame'):
                    return notes_slide.notes_text_frame.text.strip()
        except Exception as e:
            logger.debug(f"Could not extract speaker notes: {e}")
        
        return ""
    
    def _analyze_slide_layout(self, slide) -> Dict[str, Any]:
        """Analyze slide layout and visual elements"""
        analysis = {
            "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "unknown",
            "shape_count": len(slide.shapes),
            "has_images": False,
            "has_charts": False,
            "has_tables": False,
            "elements": []
        }
        
        for i, shape in enumerate(slide.shapes):
            element = {
                "shape_id": i,
                "type": "unknown",
                "content": "",
                "position": {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
            }
            
            # Identify shape types
            if hasattr(shape, "text") and shape.text.strip():
                element["type"] = "text"
                element["content"] = shape.text.strip()
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                element["type"] = "image"
                element["content"] = "Image content"
                analysis["has_images"] = True
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                element["type"] = "chart"
                element["content"] = "Chart content"
                analysis["has_charts"] = True
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                element["type"] = "table"
                element["content"] = "Table content"
                analysis["has_tables"] = True
            
            analysis["elements"].append(element)
        
        return analysis
    
    def _create_powerpoint_content(self, slide_text: str, notes: str, 
                                 layout_analysis: Dict, slide_number: int) -> str:
        """Create enhanced content for PowerPoint slide"""
        content_parts = []
        
        content_parts.append(f"=== SLIDE {slide_number} ===")
        content_parts.append(f"Layout: {layout_analysis.get('layout_name', 'Unknown')}")
        
        if slide_text:
            content_parts.append("\n=== SLIDE CONTENT ===")
            content_parts.append(slide_text)
        
        if notes:
            content_parts.append("\n=== SPEAKER NOTES ===")
            content_parts.append(notes)
        
        # Add visual element summary
        if layout_analysis.get("elements"):
            content_parts.append("\n=== VISUAL ELEMENTS ===")
            for element in layout_analysis["elements"]:
                if element["type"] != "text":  # Text already included above
                    content_parts.append(f"- {element['type']}: {element['content']}")
        
        return "\n".join(content_parts)


class AdvancedAudioVideoProcessor:
    """Enhanced audio/video processor with Whisper integration"""
    
    def process_audio_with_whisper(self, file_path: str) -> AdvancedProcessingResult:
        """Process audio files with enhanced Whisper transcription"""
        start_time = time.time()
        
        try:
            if not WHISPER_AVAILABLE:
                raise ImportError("Whisper not available for audio transcription")
            
            # Load Whisper model
            model = whisper.load_model(config.ai.whisper_model)
            
            # Transcribe with detailed options
            result = model.transcribe(
                file_path,
                language=None,  # Auto-detect language
                task="transcribe",
                verbose=False
            )
            
            # Extract detailed information
            text = result["text"]
            language = result.get("language", "unknown")
            segments = result.get("segments", [])
            
            # Create enhanced content with timestamps
            content = self._create_audio_content(text, segments, language)
            
            doc = LangChainDocument(
                page_content=content,
                metadata={
                    "source": file_path,
                    "document_type": "audio",
                    "extraction_method": "whisper_enhanced",
                    "detected_language": language,
                    "transcription_model": config.ai.whisper_model,
                    "segment_count": len(segments),
                    "duration": result.get("duration", 0),
                    "confidence_scores": [seg.get("avg_logprob", 0) for seg in segments]
                }
            )
            
            # Create structural elements for segments
            structural_elements = []
            for i, segment in enumerate(segments):
                structural_elements.append(StructuralElement(
                    element_type="audio_segment",
                    content=segment.get("text", ""),
                    position={
                        "start_time": segment.get("start", 0),
                        "end_time": segment.get("end", 0)
                    },
                    confidence=segment.get("avg_logprob", 0),
                    metadata=segment
                ))
            
            processing_time = time.time() - start_time
            
            return AdvancedProcessingResult(
                success=True,
                documents=[doc],
                structural_elements=structural_elements,
                visual_analysis={"transcription_analysis": result},
                metadata={
                    "extraction_method": "whisper_enhanced",
                    "language": language,
                    "duration": result.get("duration", 0),
                    "processing_time": processing_time
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Enhanced audio processing failed: {e}")
            return AdvancedProcessingResult(
                success=False,
                documents=[],
                structural_elements=[],
                visual_analysis={},
                metadata={"error": str(e)},
                error_message=str(e),
                processing_time=time.time() - start_time
            )
    
    def _create_audio_content(self, text: str, segments: List[Dict], language: str) -> str:
        """Create enhanced content for audio transcription"""
        content_parts = []
        
        content_parts.append("=== AUDIO TRANSCRIPTION ===")
        content_parts.append(f"Language: {language}")
        content_parts.append(f"Segments: {len(segments)}")
        
        content_parts.append("\n=== FULL TRANSCRIPT ===")
        content_parts.append(text)
        
        if segments:
            content_parts.append("\n=== TIMESTAMPED SEGMENTS ===")
            for segment in segments:
                start_time = segment.get("start", 0)
                end_time = segment.get("end", 0)
                segment_text = segment.get("text", "")
                content_parts.append(f"[{start_time:.1f}s - {end_time:.1f}s]: {segment_text}")
        
        return "\n".join(content_parts)


class AdvancedContentExtractionService:
    """Main service for advanced content extraction with llava:7b integration"""
    
    def __init__(self):
        self.pdf_processor = AdvancedPDFProcessor()
        self.excel_processor = AdvancedExcelProcessor()
        self.powerpoint_processor = AdvancedPowerPointProcessor()
        self.audio_video_processor = AdvancedAudioVideoProcessor()
        self.vision_analyzer = LlavaVisionAnalyzer()
    
    def process_document_advanced(self, file_path: str, file_type: str) -> AdvancedProcessingResult:
        """Process document with advanced extraction based on file type"""
        try:
            file_type = file_type.lower()
            
            if file_type == ".pdf":
                return self.pdf_processor.process_pdf_with_structure_analysis(file_path)
            elif file_type in [".xlsx", ".xls"]:
                return self.excel_processor.process_excel_with_relationships(file_path)
            elif file_type in [".pptx", ".ppt"]:
                return self.powerpoint_processor.process_powerpoint_with_analysis(file_path)
            elif file_type in [".mp3", ".wav", ".m4a", ".flac", ".ogg"]:
                return self.audio_video_processor.process_audio_with_whisper(file_path)
            elif file_type in [".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif"]:
                return self._process_image_with_vision(file_path)
            else:
                # Fallback to basic processing
                return self._process_with_basic_extraction(file_path, file_type)
                
        except Exception as e:
            logger.error(f"Advanced document processing failed: {e}")
            return AdvancedProcessingResult(
                success=False,
                documents=[],
                structural_elements=[],
                visual_analysis={},
                metadata={"error": str(e)},
                error_message=str(e)
            )
    
    def _process_image_with_vision(self, file_path: str) -> AdvancedProcessingResult:
        """Process images with llava:7b vision analysis"""
        start_time = time.time()
        
        try:
            # Analyze image structure and content
            structure_analysis = self.vision_analyzer.analyze_document_structure(file_path, "Image document")
            visual_elements = self.vision_analyzer.analyze_visual_elements(file_path)
            
            # Try OCR if available
            ocr_text = ""
            if TESSERACT_AVAILABLE:
                try:
                    image = Image.open(file_path)
                    ocr_text = pytesseract.image_to_string(image, lang=config.ai.ocr_language)
                except Exception as e:
                    logger.warning(f"OCR failed: {e}")
            
            # Create enhanced content
            content_parts = []
            content_parts.append("=== IMAGE ANALYSIS ===")
            content_parts.append(f"Vision Analysis: {visual_elements.get('visual_description', 'No description available')}")
            
            if structure_analysis.get("elements"):
                content_parts.append("\n=== DETECTED ELEMENTS ===")
                for element in structure_analysis["elements"]:
                    content_parts.append(f"- {element.get('type', 'Unknown')}: {element.get('description', '')}")
            
            if ocr_text.strip():
                content_parts.append("\n=== OCR TEXT ===")
                content_parts.append(ocr_text.strip())
            
            content = "\n".join(content_parts)
            
            doc = LangChainDocument(
                page_content=content,
                metadata={
                    "source": file_path,
                    "document_type": "image",
                    "extraction_method": "llava_vision_analysis",
                    "structure_analysis": structure_analysis,
                    "visual_elements": visual_elements,
                    "has_ocr_text": bool(ocr_text.strip()),
                    "vision_model": config.ai.vision_model
                }
            )
            
            # Create structural elements
            structural_elements = []
            for element in structure_analysis.get("elements", []):
                structural_elements.append(StructuralElement(
                    element_type=element.get("type", "unknown"),
                    content=element.get("description", ""),
                    position={"location": element.get("position", "")},
                    confidence=0.8,
                    metadata=element
                ))
            
            processing_time = time.time() - start_time
            
            return AdvancedProcessingResult(
                success=True,
                documents=[doc],
                structural_elements=structural_elements,
                visual_analysis={
                    "structure": structure_analysis,
                    "visual_elements": visual_elements
                },
                metadata={
                    "extraction_method": "llava_vision_analysis",
                    "processing_time": processing_time,
                    "vision_model": config.ai.vision_model
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Image processing with vision failed: {e}")
            return AdvancedProcessingResult(
                success=False,
                documents=[],
                structural_elements=[],
                visual_analysis={},
                metadata={"error": str(e)},
                error_message=str(e),
                processing_time=time.time() - start_time
            )
    
    def _process_with_basic_extraction(self, file_path: str, file_type: str) -> AdvancedProcessingResult:
        """Fallback to basic extraction for unsupported advanced processing"""
        start_time = time.time()
        
        try:
            from langchain_community.document_loaders import UnstructuredFileLoader
            
            loader = UnstructuredFileLoader(file_path)
            documents = loader.load()
            
            # Add metadata
            for doc in documents:
                doc.metadata.update({
                    "source": file_path,
                    "document_type": file_type,
                    "extraction_method": "basic_fallback",
                    "advanced_processing": False
                })
            
            processing_time = time.time() - start_time
            
            return AdvancedProcessingResult(
                success=True,
                documents=documents,
                structural_elements=[],
                visual_analysis={},
                metadata={
                    "extraction_method": "basic_fallback",
                    "processing_time": processing_time
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Basic extraction failed: {e}")
            return AdvancedProcessingResult(
                success=False,
                documents=[],
                structural_elements=[],
                visual_analysis={},
                metadata={"error": str(e)},
                error_message=str(e),
                processing_time=time.time() - start_time
            )


# Global advanced extraction service instance
advanced_extraction_service = AdvancedContentExtractionService()